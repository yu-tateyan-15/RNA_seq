#!/usr/bin/env python3
"""Build a beginner-oriented RNA-seq workflow PowerPoint deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_DIR = Path("/Users/yusuke_tateishi/Documents/RNA_seq")
PPTX_PATH = PROJECT_DIR / "Presentation_260518.pptx"

WIDE_LAYOUT = (13.333, 7.5)
COLORS = {
    "ink": RGBColor(29, 33, 41),
    "muted": RGBColor(91, 99, 112),
    "blue": RGBColor(42, 91, 215),
    "teal": RGBColor(0, 143, 142),
    "green": RGBColor(53, 148, 89),
    "orange": RGBColor(229, 126, 35),
    "red": RGBColor(195, 64, 74),
    "purple": RGBColor(111, 76, 190),
    "gray": RGBColor(236, 239, 244),
    "light_blue": RGBColor(231, 239, 255),
    "light_teal": RGBColor(226, 246, 245),
    "white": RGBColor(255, 255, 255),
}


def set_font(run, size=18, bold=False, color="ink"):
    run.font.name = "Meiryo"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(0, 0, 0)


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color="ink", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.33, 12.2, 0.55, title, size=26, bold=True)
    if subtitle:
        add_textbox(slide, 0.58, 0.92, 11.9, 0.36, subtitle, size=12, color="muted")
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["blue"]
    line.line.fill.background()


def add_card(slide, x, y, w, h, title, body, accent="blue", body_size=15):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS[accent]
    shape.line.width = Pt(1.2)
    add_textbox(slide, x + 0.14, y + 0.12, w - 0.28, 0.32, title, size=15, bold=True, color=accent)
    add_textbox(slide, x + 0.14, y + 0.48, w - 0.28, h - 0.58, body, size=body_size, color="ink")
    return shape


def add_pill(slide, x, y, w, h, text, fill="light_blue", color="blue", size=12):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.fill.background()
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=True, color=color)
    return shape


def add_bullets(slide, x, y, w, h, lines, size=15, color="ink"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for i, line in enumerate(lines):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.level = 0
        p.text = line
        p.font.name = "Meiryo"
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(4)
    return box


def add_structured_flow(slide, x=0.55, y=1.7, w=12.25):
    stages = [
        ("メタデータ整理", "品質確認", "参照辞書作成", "発現量定量", "カウント表作成", "サンプル品質確認", "差次的発現解析", "機能解釈"),
        ("サンプル表を書く\n比較を決める", "read品質を見る\n破損/adapter確認", "参照配列と注釈\nindexを作る", "readを転写産物へ\n対応づける", "転写産物を遺伝子へ\n合計する", "主成分分析/相関で\n外れ値を見る", "群間差を\n統計検定する", "機能注釈で\n意味に翻訳"),
        ("手作業\nPython", "FastQC\nMultiQC", "GENCODE\nSalmon index", "Salmon", "Python/pandas\ntx2gene", "Python\nscikit-learn", "R\nDESeq2", "R\nclusterProfiler"),
        ("samples.tsv\ncontrasts.tsv\nconfig.json", "fastqc.html\nmultiqc.html", ".fa.gz / .gtf.gz\nsalmon_index/", "quant.sf", "gene_counts.tsv", "主成分分析図\n相関図", "deseq2_*.csv\nMA plot", "Gene Ontology csv\nGSEA csv"),
    ]
    row_labels = ["解析で行うこと", "処理の意味", "使うツール", "出力ファイル"]
    row_heights = [0.46, 0.68, 0.58, 0.62]
    label_w = 1.55
    col_w = (w - label_w) / 4

    block_height = sum(row_heights)
    for block_idx in range(2):
        block_y = y + block_idx * (block_height + 0.22)
        for row_idx, row_label in enumerate(row_labels):
            yy = block_y + sum(row_heights[:row_idx])
            rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(yy), Inches(label_w), Inches(row_heights[row_idx]))
            rect.fill.solid()
            rect.fill.fore_color.rgb = COLORS["gray"]
            rect.line.color.rgb = RGBColor(210, 216, 226)
            add_textbox(slide, x + 0.08, yy + 0.1, label_w - 0.16, row_heights[row_idx] - 0.12, row_label, size=11.5, bold=True)
            for local_col in range(4):
                stage_idx = block_idx * 4 + local_col
                xx = x + label_w + local_col * col_w
                cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(xx), Inches(yy), Inches(col_w), Inches(row_heights[row_idx]))
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["light_blue"] if row_idx == 0 else COLORS["white"]
                cell.line.color.rgb = RGBColor(210, 216, 226)
                if row_idx == 0:
                    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(xx + 0.08), Inches(yy + 0.09), Inches(0.24), Inches(0.24))
                    circle.fill.solid()
                    circle.fill.fore_color.rgb = COLORS["white"]
                    circle.line.color.rgb = COLORS["blue"]
                    add_textbox(slide, xx + 0.08, yy + 0.105, 0.24, 0.15, str(stage_idx + 1), size=8.8, bold=True, align=PP_ALIGN.CENTER)
                    text_x = xx + 0.36
                    text_w = col_w - 0.42
                else:
                    text_x = xx + 0.06
                    text_w = col_w - 0.12
                add_textbox(
                    slide,
                    text_x,
                    yy + 0.06,
                    text_w,
                    row_heights[row_idx] - 0.08,
                    stages[row_idx][stage_idx],
                    size=11.2 if row_idx else 12.0,
                    bold=row_idx == 0,
                    align=PP_ALIGN.CENTER,
                )


def add_line_icon(slide, kind, x, y, size=0.9, accent="blue"):
    left = Inches(x)
    top = Inches(y)
    s = Inches(size)
    color = COLORS[accent]
    if kind == "document":
        page = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, s * 0.72, s * 0.9)
        page.fill.solid()
        page.fill.fore_color.rgb = COLORS["white"]
        page.line.color.rgb = color
        for i in range(3):
            line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + s * 0.12, top + s * (0.22 + i * 0.18), s * 0.46, s * 0.025)
            line.fill.solid()
            line.fill.fore_color.rgb = color
            line.line.fill.background()
    elif kind == "magnifier":
        lens = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, s * 0.58, s * 0.58)
        lens.fill.background()
        lens.line.color.rgb = color
        lens.line.width = Pt(2)
        handle = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + s * 0.48, top + s * 0.48, left + s * 0.82, top + s * 0.82)
        handle.line.color.rgb = color
        handle.line.width = Pt(2)
    elif kind == "database":
        for i in range(3):
            oval = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top + s * (0.12 + i * 0.22), s * 0.8, s * 0.22)
            oval.fill.background()
            oval.line.color.rgb = color
            oval.line.width = Pt(1.4)
        body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top + s * 0.23, s * 0.8, s * 0.45)
        body.fill.background()
        body.line.color.rgb = color
    elif kind == "grid":
        for i in range(4):
            vx = left + s * i * 0.2
            vline = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, vx, top, vx, top + s * 0.78)
            vline.line.color.rgb = color
            hline = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top + s * i * 0.2, left + s * 0.78, top + s * i * 0.2)
            hline.line.color.rgb = color
    elif kind == "chart":
        axis1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top + s * 0.78, left + s * 0.88, top + s * 0.78)
        axis2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top + s * 0.78, left, top)
        axis1.line.color.rgb = color
        axis2.line.color.rgb = color
        pts = [(0.12, 0.62), (0.28, 0.45), (0.46, 0.52), (0.68, 0.2)]
        for a, b in zip(pts, pts[1:]):
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + s * a[0], top + s * a[1], left + s * b[0], top + s * b[1])
            line.line.color.rgb = color
            line.line.width = Pt(2)
    elif kind == "network":
        coords = [(0.15, 0.25), (0.55, 0.16), (0.35, 0.58), (0.72, 0.62)]
        for a, b in [(0, 1), (0, 2), (1, 2), (2, 3)]:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + s * coords[a][0], top + s * coords[a][1], left + s * coords[b][0], top + s * coords[b][1])
            line.line.color.rgb = color
        for cx, cy in coords:
            node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + s * cx - s * 0.055, top + s * cy - s * 0.055, s * 0.11, s * 0.11)
            node.fill.solid()
            node.fill.fore_color.rgb = COLORS["white"]
            node.line.color.rgb = color


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(WIDE_LAYOUT[0])
    prs.slide_height = Inches(WIDE_LAYOUT[1])

    # 1
    slide = blank_slide(prs)
    add_textbox(slide, 0.75, 1.0, 11.6, 0.7, "RNA sequencing解析の全体フロー", size=34, bold=True)
    add_textbox(slide, 0.8, 1.82, 10.8, 0.45, "塩基配列readファイルから生物学的解釈までを、ファイル・ツール・数式で理解する", size=18, color="muted")
    add_structured_flow(slide, x=0.55, y=2.0, w=12.25)

    # 2
    slide = blank_slide(prs)
    add_title(slide, "全体像: RNA sequencingはファイルを段階的に変換する解析", "各ステップは前の出力に依存するため、順番を飛ばすと下流の解釈が崩れる")
    add_structured_flow(slide, x=0.55, y=1.42, w=12.25)
    add_textbox(slide, 0.75, 6.55, 11.7, 0.35, "読み方: 上段で解析内容、下段で使うツールと出力ファイルを対応づける。", 16, True)

    # 3
    slide = blank_slide(prs)
    add_title(slide, "拡張子で見るRNA sequencingの登場人物", "拡張子を見ると、そのファイルが解析のどの段階にいるか分かる")
    rows = [
        (".fq.gz", "圧縮FASTQ（塩基配列read）", "生の塩基配列read。解析の出発点"),
        (".fa.gz", "圧縮FASTA（参照配列）", "参照転写産物配列。Salmon indexの材料"),
        (".gtf.gz", "圧縮GTF（遺伝子注釈）", "遺伝子/転写産物/exonの注釈"),
        (".tsv", "タブ区切り表", "samples.tsv, tx2gene.tsv, gene_counts.tsv"),
        (".sf", "Salmon出力", "転写産物ごとのTPM（長さ補正発現量）/推定read数"),
        (".csv", "カンマ区切り表", "差次的発現解析結果やDEG（差次的発現遺伝子）リスト"),
        (".html", "レポート", "MultiQCやFastQCの結果をブラウザで見る"),
        (".png/.pdf", "図", "主成分分析図、heatmap、volcano plot、MA plot"),
    ]
    x = [0.75, 2.45, 5.0]
    add_textbox(slide, x[0], 1.48, 1.4, 0.3, "拡張子", 13, True, "blue")
    add_textbox(slide, x[1], 1.48, 2.2, 0.3, "意味", 13, True, "blue")
    add_textbox(slide, x[2], 1.48, 7.4, 0.3, "この解析での役割", 13, True, "blue")
    for i, row in enumerate(rows):
        y = 1.9 + i * 0.53
        fill = "gray" if i % 2 == 0 else "white"
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(y - 0.06), Inches(12.1), Inches(0.42))
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS[fill]
        rect.line.fill.background()
        add_textbox(slide, x[0], y, 1.5, 0.25, row[0], 12, True)
        add_textbox(slide, x[1], y, 2.2, 0.25, row[1], 12)
        add_textbox(slide, x[2], y, 7.4, 0.25, row[2], 12)
    add_textbox(slide, 0.75, 6.6, 11.4, 0.38, "拡張子は「データの状態」を表す。FASTQはまだ生read、CSV/TSVは解析済みの表。", 14, True, "teal")

    # 4
    slide = blank_slide(prs)
    add_title(slide, "Step 0: メタデータと比較設計", "データ処理より前に、何を比較するかを固定する")
    add_line_icon(slide, "document", 11.4, 1.55, size=0.95, accent="blue")
    add_card(slide, 0.7, 1.55, 3.7, 2.2, "入力", "raw_data/*.fq.gz\n実験メモ\nサンプル名・処理条件", "blue")
    add_card(slide, 4.75, 1.55, 3.7, 2.2, "出力", "metadata/samples.tsv\nmetadata/contrasts.tsv\nconfig/analysis_config.json", "teal")
    add_card(slide, 8.8, 1.55, 3.7, 2.2, "なぜ必要か", "DESeq2はmetadataのconditionを使う。\n設計が曖昧だと、何を比べた統計か分からない。", "orange")
    add_textbox(slide, 0.75, 4.15, 11.7, 0.45, "例: contrast = NAC_S2_2h vs Non は「NAC_S2_2hがNonより高い/低い」を検定する", 16, bold=True)
    add_textbox(slide, 0.78, 5.05, 11.5, 0.9, "数式的には、各sample sにcondition c_s を割り当て、発現量を condition で説明する準備をしている。\nDESeq2の設計式: y_gs ~ condition_s", 15)

    # 5
    slide = blank_slide(prs)
    add_title(slide, "Step 1: 生readの品質確認", "塩基配列readファイルを定量へ進めてよいかを確認する入口検査")
    add_line_icon(slide, "magnifier", 11.35, 1.55, size=1.0, accent="teal")
    add_card(slide, 0.7, 1.55, 3.6, 2.55, "使うツール", "FastQC: FASTQごとの品質\nMultiQC: 複数FastQCを統合\ngzip check: 圧縮ファイル破損確認", "blue")
    add_card(slide, 4.85, 1.55, 3.6, 2.55, "見る指標", "read quality（塩基品質）\nadapter contamination（人工配列混入）\nGC content（G/C塩基比率）\nread数\n重複率", "teal")
    add_card(slide, 9.0, 1.55, 3.3, 2.55, "なぜ先にやるか", "壊れた塩基配列readファイルや強いadapter汚染を使うと、後のmapping/定量が歪む。", "orange")
    add_textbox(slide, 0.8, 4.55, 5.8, 0.32, "read数の基本式", 14, True, "blue")
    add_textbox(slide, 0.8, 4.94, 5.8, 0.5, "read数 = FASTQの行数 / 4", 20, True)
    add_textbox(slide, 7.0, 4.55, 5.3, 0.32, "GC content（G/C塩基比率）", 14, True, "blue")
    add_textbox(slide, 7.0, 4.94, 5.3, 0.8, "GC% = (G + C) / (A + T + G + C) × 100", 18, True)
    add_textbox(slide, 0.8, 6.25, 11.5, 0.32, "今回: 18個のFASTQはすべて存在、gzip OK、GC contentは48-49%、adapter contentはpass。", 14, True, "green")

    # 6
    slide = blank_slide(prs)
    add_title(slide, "Step 2: 参照ファイルとSalmon index", "塩基配列readを遺伝子・転写産物へ対応づけるための辞書を作る")
    add_line_icon(slide, "database", 11.35, 1.55, size=1.0, accent="purple")
    add_card(slide, 0.7, 1.55, 3.5, 2.5, ".fa.gz", "GENCODE transcript FASTA。\nヒト転写産物配列の一覧。\nSalmon indexの材料。", "blue")
    add_card(slide, 4.45, 1.55, 3.5, 2.5, ".gtf.gz", "GENCODE annotation GTF。\n遺伝子/転写産物/exonの対応と座標。\ntx2gene作成に使う。", "teal")
    add_card(slide, 8.2, 1.55, 3.9, 2.5, "salmon_index/", "FASTAをSalmonが高速検索できる内部形式に変換したもの。\n人間が読む表ではない。", "purple")
    add_textbox(slide, 0.85, 4.45, 11.5, 0.55, "転写産物から遺伝子への集計", 16, True, "blue")
    add_textbox(slide, 0.85, 5.0, 11.5, 0.7, "gene_count(g, s) = Σ NumReads(t, s)  for transcripts t belonging to gene g", 20, True)
    add_textbox(slide, 0.85, 6.05, 11.5, 0.35, "tx2gene.tsv が「転写産物 t は遺伝子 g に属する」という対応表。", 14)

    # 7
    slide = blank_slide(prs)
    add_title(slide, "Step 3: Salmonによる発現量定量", "塩基配列readを転写産物ごとの発現量に変換する")
    add_line_icon(slide, "chart", 11.35, 1.55, size=1.0, accent="green")
    add_card(slide, 0.75, 1.5, 3.6, 2.4, "入力", "FASTQ（塩基配列read）: .fq.gz\nSalmon index\nlibrary type: A", "blue")
    add_card(slide, 4.85, 1.5, 3.6, 2.4, "出力", "quant.sf\nTPM（長さ補正発現量）\nNumReads（推定read数）\nmapping rate", "teal")
    add_card(slide, 8.95, 1.5, 3.4, 2.4, "なぜ必要か", "FASTQはreadの羅列。\n統計解析には遺伝子ごとの数値表が必要。", "orange")
    add_textbox(slide, 0.85, 4.35, 11.5, 0.5, "mapping rate", 16, True, "blue")
    add_textbox(slide, 0.85, 4.9, 11.5, 0.55, "mapping_rate = mapped_reads / processed_reads × 100", 20, True)
    add_textbox(slide, 0.85, 5.82, 11.5, 0.42, "今回例: 11,667,534 / 12,939,243 × 100 = 90.17%。入口として良い値。", 14, True, "green")

    # 8
    slide = blank_slide(prs)
    add_title(slide, "Step 4: Count matrixとサンプル品質確認", "統計前にサンプル全体の構造を見る")
    add_line_icon(slide, "grid", 11.35, 1.55, size=1.0, accent="blue")
    add_card(slide, 0.7, 1.5, 3.7, 2.35, "gene_counts.tsv", "行 = 遺伝子\n列 = サンプル\n値 = 遺伝子count", "blue")
    add_card(slide, 4.8, 1.5, 3.7, 2.35, "見る図", "library size（総count）\ndetected genes（検出遺伝子数）\n主成分分析\n相関heatmap\n階層的クラスタリング", "teal")
    add_card(slide, 8.9, 1.5, 3.5, 2.35, "なぜ必要か", "外れ値やbatchがあると、DESeq2の結果を誤解する。", "orange")
    add_textbox(slide, 0.8, 4.3, 5.8, 0.4, "library size（総count）", 14, True, "blue")
    add_textbox(slide, 0.8, 4.72, 5.8, 0.5, "library_size_s = Σ_g count_gs", 19, True)
    add_textbox(slide, 6.8, 4.3, 5.5, 0.4, "logCPM（log変換したCPM）", 14, True, "blue")
    add_textbox(slide, 6.8, 4.72, 5.5, 0.5, "logCPM_gs = log2(CPM_gs + 1)", 19, True)
    add_textbox(slide, 0.8, 6.05, 11.5, 0.35, "今回: library sizeと検出遺伝子数に大きな破綻なし。NAC_S2_2hは主成分分析/clusterで別にまとまる。", 14, True, "green")

    # 9
    slide = blank_slide(prs)
    add_title(slide, "Step 5: DESeq2で差次的発現解析", "count matrixを、遺伝子ごとの統計結果に変換する")
    add_line_icon(slide, "chart", 11.35, 1.55, size=1.0, accent="orange")
    add_card(slide, 0.75, 1.5, 3.8, 2.45, "入力", "gene_counts.tsv\nsamples.tsv\ncontrasts.tsv\ndesign: ~ condition", "blue")
    add_card(slide, 4.9, 1.5, 3.8, 2.45, "出力", "log2FoldChange\npvalue\npadj\nDEG list（差次的発現遺伝子リスト）\nMA plot", "teal")
    add_card(slide, 9.05, 1.5, 3.3, 2.45, "なぜ必要か", "平均差だけでなく、replicate間のばらつきと多重検定を考慮する。", "orange")
    add_textbox(slide, 0.85, 4.4, 11.4, 0.52, "DESeq2の基本モデル", 16, True, "blue")
    add_textbox(slide, 0.85, 4.9, 11.4, 0.62, "count_gs ~ NegativeBinomial(mean_gs, dispersion_g)", 20, True)
    add_textbox(slide, 0.85, 5.75, 11.4, 0.45, "log2FC = log2(mean_test / mean_reference)。padjは多重検定補正後のp値。", 15)
    add_textbox(slide, 0.85, 6.35, 11.4, 0.35, "今回: NAC_S2_2h vs Nonで55個のDEG、NAC_S2_2h vs Oxi_2hで54個のDEG。", 14, True, "green")

    # 10
    slide = blank_slide(prs)
    add_title(slide, "Step 6: 生物学的解釈", "DEGリスト（差次的発現遺伝子リスト）をpathwayや機能の言葉に翻訳する")
    add_line_icon(slide, "network", 11.35, 1.55, size=1.0, accent="purple")
    add_card(slide, 0.75, 1.5, 3.8, 2.35, "使うDB/ツール", "clusterProfiler\nGene Ontology Biological Process\nGene Set Enrichment Analysis\nReactome/UniProtで確認", "blue")
    add_card(slide, 4.9, 1.5, 3.8, 2.35, "見る出力", "volcano plot\nDEG heatmap\nGene Ontology enrichment\nGSEA NES（正規化enrichment score）", "teal")
    add_card(slide, 9.05, 1.5, 3.3, 2.35, "なぜ必要か", "DEG名の羅列だけでは生物学的ストーリーにならない。", "orange")
    add_textbox(slide, 0.85, 4.35, 11.4, 0.5, "DEG（差次的発現遺伝子）集合 D と機能term集合 T の重なり", 16, True, "blue")
    add_textbox(slide, 0.85, 4.9, 11.4, 0.6, "k = |D ∩ T|", 25, True)
    add_textbox(slide, 0.85, 5.75, 11.4, 0.5, "k が偶然にしては大きいかを検定する。今回、NAC_S2_2hで酸化的リン酸化・ミトコンドリア呼吸鎖が強く出た。", 14, True, "green")

    # 11
    slide = blank_slide(prs)
    add_title(slide, "なぜこの順番なのか", "下流解析は上流の品質と変換に依存する")
    add_structured_flow(slide, x=0.55, y=1.42, w=12.25)
    add_textbox(slide, 0.85, 6.55, 11.4, 0.38, "各段階の目的は「次の段階に渡してよい入力か」を保証すること。", 16, True)

    # 12
    slide = blank_slide(prs)
    add_title(slide, "このプロジェクトで実際に見る成果物", "どのファイルを開けば、どの段階の結果が分かるか")
    outputs = [
        ("品質確認", "results/qc/multiqc_report.html", "FASTQ品質、adapter、GC content、read数"),
        ("定量", "results/quant/salmon/*/quant.sf", "転写産物ごとのTPM（長さ補正発現量）/NumReads"),
        ("count", "results/counts/gene_counts.tsv", "DESeq2に渡すgene × sample表"),
        ("サンプル品質確認", "results/sample_qc/*.png", "主成分分析、相関、library size"),
        ("差次的発現解析", "results/de/deseq2_*.csv", "log2FC、padj、DEG（差次的発現遺伝子）候補"),
        ("解釈", "results/report/*.png / results/enrichment/*", "volcano、heatmap、Gene Ontology/GSEA"),
    ]
    for i, (stage, path, meaning) in enumerate(outputs):
        y = 1.55 + i * 0.76
        add_pill(slide, 0.75, y, 1.4, 0.38, stage, fill="light_blue", color="blue")
        add_textbox(slide, 2.35, y + 0.04, 4.2, 0.3, path, 12, True)
        add_textbox(slide, 6.75, y + 0.04, 5.6, 0.3, meaning, 12)
    add_textbox(slide, 0.8, 6.55, 11.3, 0.35, "解析の読み方: まず品質確認 → サンプル品質確認 → 差次的発現解析 → enrichmentの順に信頼性を確認する。", 14, True, "teal")

    prs.save(PPTX_PATH)
    print(f"wrote {PPTX_PATH}")


if __name__ == "__main__":
    build_deck()
